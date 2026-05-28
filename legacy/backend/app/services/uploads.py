from __future__ import annotations

import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path

from fastapi import HTTPException, UploadFile, status
from pypdf import PdfReader
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.db.models import DocumentChangeRequest, Role, UploadDocument, UploadStatus, User
from app.services.audit import log_action
from app.services.vector_index import safe_delete_upload_from_session, safe_index_upload_text, safe_promote_upload_to_base


def _safe_filename(name: str) -> str:
    return ''.join(ch for ch in name if ch.isalnum() or ch in ('-', '_', '.', ' ')).strip() or 'upload.bin'


def _extract_pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = [page.extract_text() or '' for page in reader.pages]
    return '\n'.join(parts).strip()


def _extract_docx_text(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return '\n'.join(p.text for p in doc.paragraphs).strip()


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                texts.append(shape.text)
    return '\n'.join(texts).strip()


def extract_text(path: Path, mime_type: str) -> str:
    suffix = path.suffix.lower()
    if suffix in {'.txt', '.md'}:
        return path.read_text(encoding='utf-8', errors='ignore').strip()
    if suffix == '.pdf':
        return _extract_pdf_text(path)
    if suffix == '.docx':
        return _extract_docx_text(path)
    if suffix == '.pptx':
        return _extract_pptx_text(path)
    if suffix in {'.jpg', '.jpeg', '.png'}:
        return '图像文档（待 Qwen-VL 解析）'
    return ''


def get_upload_by_id(db: Session, upload_id: str, include_deleted: bool = False) -> UploadDocument:
    query = db.query(UploadDocument).filter(UploadDocument.id == upload_id)
    if not include_deleted:
        query = query.filter(UploadDocument.is_deleted.is_(False))
    row = query.first()
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail='Upload not found')
    return row


def create_upload(db: Session, user: User, session_id: str, file: UploadFile) -> UploadDocument:
    settings = get_settings()
    if not session_id:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail='session_id is required')

    safe_name = _safe_filename(file.filename or 'upload.bin')
    upload_id = str(uuid.uuid4())

    user_dir = settings.upload_root / user.id / session_id / upload_id
    user_dir.mkdir(parents=True, exist_ok=True)
    storage_path = user_dir / safe_name

    with storage_path.open('wb') as out:
        shutil.copyfileobj(file.file, out)

    size_bytes = storage_path.stat().st_size
    extracted = extract_text(storage_path, file.content_type or 'application/octet-stream')

    row = UploadDocument(
        id=upload_id,
        owner_id=user.id,
        session_id=session_id,
        original_filename=file.filename or safe_name,
        stored_filename=safe_name,
        mime_type=file.content_type or 'application/octet-stream',
        size_bytes=size_bytes,
        storage_path=str(storage_path),
        extracted_text=extracted,
        status=UploadStatus.draft,
        is_deleted=False,
    )
    db.add(row)
    db.commit()
    db.refresh(row)

    safe_index_upload_text(
        session_id=session_id,
        upload_id=row.id,
        text=extracted,
        payload={'owner_id': user.id, 'status': row.status.value, 'source_domain': 'upload'},
    )

    log_action(
        db,
        action='upload.create',
        target_type='upload',
        target_id=row.id,
        actor=user,
        details={'session_id': session_id, 'filename': row.original_filename},
    )

    return row


def list_uploads_for_session(db: Session, user: User, session_id: str) -> list[UploadDocument]:
    query = db.query(UploadDocument).filter(
        UploadDocument.session_id == session_id,
        UploadDocument.is_deleted.is_(False),
    )
    if user.role != Role.admin:
        query = query.filter(UploadDocument.owner_id == user.id)
    return query.order_by(UploadDocument.created_at.desc()).all()


def submit_upload_for_review(db: Session, user: User, upload_id: str, note: str | None = None) -> UploadDocument:
    row = get_upload_by_id(db, upload_id)
    if user.role != Role.admin and row.owner_id != user.id:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail='No permission to submit this upload')
    if row.status != UploadStatus.draft:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail='Only draft uploads can be submitted')

    row.status = UploadStatus.pending_review
    row.submitted_at = datetime.now(timezone.utc)
    if note:
        row.review_reason = note
    db.commit()
    db.refresh(row)

    log_action(
        db,
        action='upload.submit_for_review',
        target_type='upload',
        target_id=row.id,
        actor=user,
        details={'session_id': row.session_id, 'note': note},
    )
    return row


def delete_draft_upload(db: Session, user: User, upload_id: str) -> None:
    row = get_upload_by_id(db, upload_id)
    if row.owner_id != user.id and user.role != Role.admin:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail='No permission to delete this upload')
    if row.status != UploadStatus.draft:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail='Only draft uploads can be deleted')

    path = Path(row.storage_path)
    if path.exists():
        try:
            path.unlink()
            parent = path.parent
            if parent.exists() and not any(parent.iterdir()):
                parent.rmdir()
        except OSError:
            pass

    safe_delete_upload_from_session(row.session_id, row.id)

    db.delete(row)
    db.commit()

    log_action(db, action='upload.delete_draft', target_type='upload', target_id=upload_id, actor=user)


def soft_delete_upload(db: Session, admin: User, upload_id: str) -> UploadDocument:
    row = get_upload_by_id(db, upload_id, include_deleted=True)
    if row.is_deleted:
        return row

    row.is_deleted = True
    row.deleted_at = datetime.now(timezone.utc)
    row.deleted_by = admin.id
    db.commit()
    db.refresh(row)

    safe_delete_upload_from_session(row.session_id, row.id)

    log_action(
        db,
        action='admin.document.soft_delete',
        target_type='upload',
        target_id=row.id,
        actor=admin,
    )
    return row


def approve_submission(db: Session, admin: User, upload_id: str, reason: str | None = None) -> UploadDocument:
    row = get_upload_by_id(db, upload_id)
    if row.status != UploadStatus.pending_review:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail='Submission is not pending review')

    row.status = UploadStatus.approved
    row.reviewed_by = admin.id
    row.review_reason = reason
    row.reviewed_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(row)

    safe_promote_upload_to_base(
        upload_id=row.id,
        text=row.extracted_text or '',
        payload={
            'source_domain': 'base',
            'owner_id': row.owner_id,
            'filename': row.original_filename,
            'approved_by': admin.id,
        },
    )

    log_action(
        db,
        action='admin.review.approve',
        target_type='upload',
        target_id=row.id,
        actor=admin,
        details={'reason': reason},
    )

    return row


def reject_submission(db: Session, admin: User, upload_id: str, reason: str | None = None) -> UploadDocument:
    row = get_upload_by_id(db, upload_id)
    if row.status != UploadStatus.pending_review:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail='Submission is not pending review')

    row.status = UploadStatus.rejected
    row.reviewed_by = admin.id
    row.review_reason = reason
    row.reviewed_at = datetime.now(timezone.utc)
    db.commit()
    db.refresh(row)

    log_action(
        db,
        action='admin.review.reject',
        target_type='upload',
        target_id=row.id,
        actor=admin,
        details={'reason': reason},
    )

    return row


def get_submission(db: Session, admin: User, upload_id: str) -> UploadDocument:
    row = get_upload_by_id(db, upload_id)
    if row.status != UploadStatus.pending_review:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail='Submission not found')
    return row


def list_pending_submissions(db: Session) -> list[UploadDocument]:
    return (
        db.query(UploadDocument)
        .filter(
            UploadDocument.status == UploadStatus.pending_review,
            UploadDocument.is_deleted.is_(False),
        )
        .order_by(UploadDocument.submitted_at.desc())
        .all()
    )


def create_change_request(
    db: Session,
    user: User,
    upload_id: str,
    proposed_filename: str | None,
    proposed_extracted_text: str | None,
    reason: str | None,
) -> DocumentChangeRequest:
    row = get_upload_by_id(db, upload_id)
    if user.role != Role.admin and row.owner_id != user.id:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail='No permission to request this change')

    request = DocumentChangeRequest(
        document_id=row.id,
        requester_id=user.id,
        proposed_filename=proposed_filename,
        proposed_extracted_text=proposed_extracted_text,
        reason=reason,
        status='pending',
    )
    db.add(request)
    db.commit()
    db.refresh(request)
    log_action(
        db,
        action='document.change_request.create',
        target_type='document_change_request',
        target_id=str(request.id),
        actor=user,
        details={'document_id': row.id},
    )
    return request

